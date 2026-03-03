#!/usr/bin/env python3
"""
Minimal example: McKinsey-style PPT with Cover + Content + Source slides.
Uses the design system defined in SKILL.md.
"""

import os
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Color Palette ──
NAVY      = RGBColor(0x05, 0x1C, 0x2C)
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY  = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)

# ── Font Sizes ──
TITLE_SIZE      = Pt(22)
BODY_SIZE       = Pt(14)
SUB_HEADER_SIZE = Pt(18)
HEADER_SIZE     = Pt(28)
SMALL_SIZE      = Pt(9)

# ── Helper Functions ──

def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)

def add_text(slide, left, top, width, height, text, font_size=Pt(14),
             font_name='Arial', font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = font_name
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for run in p.runs:
        set_ea_font(run, ea_font)
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=BLACK, width=Pt(0.5)):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = width
    sp = c._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)
    return c

def add_action_title(slide, text, title_size=Pt(22)):
    add_text(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9),
             text, font_size=title_size, font_color=BLACK, bold=True,
             font_name='Georgia', ea_font='KaiTi', anchor=MSO_ANCHOR.MIDDLE)
    add_line(slide, Inches(0.8), Inches(1.05), Inches(12.5), Inches(1.05),
             color=BLACK, width=Pt(0.5))

def add_source(slide, text, y=Inches(7.05)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(0.3),
             text, font_size=Pt(9), font_color=MED_GRAY)

def cleanup_theme(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if 'theme' in item.filename.lower() and item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    for eff_list in root.findall('.//a:effectStyleLst/a:effectStyle/a:effectLst', ns):
                        for shadow in eff_list.findall('a:outerShdw', ns):
                            eff_list.remove(shadow)
                        for shadow in eff_list.findall('a:innerShdw', ns):
                            eff_list.remove(shadow)
                    for eff_style in root.findall('.//a:effectStyleLst/a:effectStyle', ns):
                        for s3d in eff_style.findall('a:scene3d', ns):
                            eff_style.remove(s3d)
                        for sp3d in eff_style.findall('a:sp3d', ns):
                            eff_style.remove(sp3d)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)

# ── Build Presentation ──

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Cover
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, prs.slide_width, Inches(0.05), NAVY)
    add_text(s1, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
             'Sample Presentation', font_size=Pt(44), font_name='Georgia',
             font_color=NAVY, bold=True)
    add_text(s1, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             'McKinsey-style Design System Demo', font_size=Pt(24),
             font_color=DARK_GRAY)
    add_text(s1, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             'Minimal Example  |  2026', font_size=BODY_SIZE, font_color=MED_GRAY)
    add_line(s1, Inches(1), Inches(6.8), Inches(4), Inches(6.8),
             color=NAVY, width=Pt(2))

    # Slide 2: Content
    s2 = prs.slides.add_slide(blank)
    add_action_title(s2, 'Key Findings Overview')
    items = [
        'Clean typography hierarchy ensures readability',
        'Flat design with no shadows or 3D effects',
        'Consistent color palette across all slides',
        'Proper East Asian font handling for Chinese text',
    ]
    for i, item in enumerate(items):
        y = Inches(1.6) + Inches(0.6) * i
        add_text(s2, Inches(1.2), y, Inches(10), Inches(0.5), item)
        if i < len(items) - 1:
            add_line(s2, Inches(1.2), y + Inches(0.55),
                     Inches(11.2), y + Inches(0.55), color=LINE_GRAY)
    add_source(s2, 'Source: Mck-ppt-design-skill v1.0.0')

    # Save & cleanup
    outpath = 'minimal_output.pptx'
    prs.save(outpath)
    cleanup_theme(outpath)
    print(f'Created: {outpath} ({os.path.getsize(outpath):,} bytes)')

if __name__ == '__main__':
    main()
